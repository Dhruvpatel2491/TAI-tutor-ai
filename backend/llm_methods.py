import os
os.environ["TF_ENABLE_ONEDNN_OPTS"] = "0"

import argparse
import logging
import json
import shutil
from pathlib import Path

import fitz                # PyMuPDF
from pptx import Presentation
from tree_sitter import Language, Parser
import nbformat
from pydantic import BaseModel
import ast

# llama_index imports (keep matching your environment)
from llama_index.core import (
    SimpleDirectoryReader,
    VectorStoreIndex,
    StorageContext,
    load_index_from_storage,
)
from llama_index.core.node_parser import SimpleNodeParser
from llama_index.llms.ollama import Ollama
from llama_index.embeddings.ollama import OllamaEmbedding
from llama_index.core.settings import Settings

# ----------------------------- CONFIG -----------------------------
INDEX_DIR = "./index_store"
DATA_DIR = "./trial-data/                "

# If you want to require a specific file for the persisted index, add it here.
# LlamaIndex currently expects a `docstore.json` for simple_docstore, so include it.
EXPECTED_PERSIST_FILES = ["docstore.json"]

# Add embedding metadata filename
EMBEDDING_META_FILENAME = "embeddings_meta.json"

# ----------------------------- LOGGING -----------------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

# ----------------------------- PARSERS -----------------------------


def extract_text_from_pdf(filepath: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    try:
        with fitz.open(filepath) as doc:
            text = [page.get_text("text") for page in doc]
        return "\n".join(text)
    except Exception as e:
        logger.error(f"PDF extraction failed for {filepath}: {e}")
        return ""


def extract_text_from_pptx(filepath: str) -> str:
    """Extract slide-level text using python-pptx."""
    try:
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())
            if content:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(content))
        return "\n\n".join(slides_text)
    except Exception as e:
        logger.error(f"PPTX extraction failed for {filepath}: {e}")
        return ""


def extract_code_with_treesitter(filepath: str, language: str = "python") -> str:
    """Parse source code using Tree-sitter (multi-language)."""
    lang_so = "build/my-languages.so"
    # Build once if missing
    if not os.path.exists(lang_so):
        logger.info("Building Tree-sitter language library (this may take a moment)...")
        os.makedirs("vendor", exist_ok=True)
        # Ensure the vendor directories exist and contain tree-sitter grammars.
        # The following list is what the original script expected; adjust if you don't have all repos.
        Language.build_library(
            lang_so,
            [
                "vendor/tree-sitter-python",
                "vendor/tree-sitter-javascript",
                "vendor/tree-sitter-cpp",
                "vendor/tree-sitter-java",
            ],
        )

    # If requested language not included in the built .so, Language() may raise.
    try:
        LANGUAGE = Language(lang_so, language)
    except Exception as e:
        logger.error(f"Tree-sitter Language init failed for '{language}': {e}. Falling back to raw code.")
        return Path(filepath).read_text(encoding="utf-8", errors="ignore")

    parser = Parser()
    parser.set_language(LANGUAGE)

    code = Path(filepath).read_text(encoding="utf-8", errors="ignore")
    # parse to ensure tree creation (we don't use the tree here, but it's available)
    try:
        _ = parser.parse(bytes(code, "utf8"))
    except Exception as e:
        logger.warning(f"Tree-sitter parse warning for {filepath}: {e}")
    return code


class NotebookCell(BaseModel):
    cell_type: str
    source: str
    metadata: dict


def extract_notebook_cells(filepath: str) -> str:
    """Extract notebook code and markdown cells as JSON."""
    try:
        nb = nbformat.read(filepath, as_version=4)
        cells = [
            NotebookCell(
                cell_type=c.cell_type,
                source="".join(c.source) if isinstance(c.source, (list, tuple)) else str(c.source),
                metadata=getattr(c, "metadata", {}),
            ).dict()
            for c in nb.cells
        ]
        return json.dumps(cells, indent=2)
    except Exception as e:
        logger.error(f"Notebook parse failed for {filepath}: {e}")
        return ""


def extract_code_with_ast(filepath: str) -> str:
    """Analyze Python code structure with AST (functions/classes)."""
    try:
        src = Path(filepath).read_text(encoding="utf-8")
        tree = ast.parse(src)
        funcs = [n.name for n in ast.walk(tree) if isinstance(n, ast.FunctionDef)]
        classes = [n.name for n in ast.walk(tree) if isinstance(n, ast.ClassDef)]
        summary = f"Functions: {funcs}\nClasses: {classes}\n\n{src}"
        return summary
    except Exception as e:
        logger.error(f"AST parse failed for {filepath}: {e}")
        return ""


def load_multimodal_documents(directory: str):
    """Load PDFs, PPTXs, Notebooks, and code intelligently."""
    from llama_index.core.schema import Document
    docs = []

    if not os.path.exists(directory):
        logger.warning(f"Data directory does not exist: {directory}")
        return docs

    for root, _, files in os.walk(directory):
        for file in files:
            path = os.path.join(root, file)
            ext = Path(file).suffix.lower()

            try:
                if ext == ".pdf":
                    logger.info(f"Parsing PDF: {path}")
                    text = extract_text_from_pdf(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "pdf", "path": path}))

                elif ext == ".pptx":
                    logger.info(f"Parsing PPTX: {path}")
                    text = extract_text_from_pptx(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "pptx", "path": path}))

                elif ext == ".ipynb":
                    logger.info(f"Parsing notebook: {path}")
                    text = extract_notebook_cells(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "notebook", "path": path}))

                elif ext == ".py":
                    logger.info(f"Parsing Python code (AST): {path}")
                    text = extract_code_with_ast(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "python", "path": path}))

                elif ext in [".java", ".cpp", ".js", ".c"]:
                    logger.info(f"Parsing {ext} code with Tree-sitter: {path}")
                    lang = ext.strip(".")
                    text = extract_code_with_treesitter(path, language=lang)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "code", "lang": lang, "path": path}))

                else:
                    logger.info(f"Fallback reader for: {path}")
                    fallback = SimpleDirectoryReader(input_files=[path])
                    loaded = fallback.load_data()
                    docs.extend(loaded)

            except Exception as e:
                logger.error(f"Failed to parse {path}: {e}")

    return docs

# New helper: load docs for an explicit list of file paths (absolute paths)
def load_documents_for_paths(paths):
    """Load a list of file paths into llama_index Documents (uses the same extractors)."""
    from llama_index.core.schema import Document
    docs = []
    for path in paths:
        try:
            if not os.path.exists(path):
                logger.warning(f"Path not found while building selective docs: {path}")
                continue
            ext = Path(path).suffix.lower()
            if ext == ".pdf":
                logger.info(f"[delta] Parsing PDF: {path}")
                text = extract_text_from_pdf(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "pdf", "path": path}))
            elif ext == ".pptx":
                logger.info(f"[delta] Parsing PPTX: {path}")
                text = extract_text_from_pptx(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "pptx", "path": path}))
            elif ext == ".ipynb":
                logger.info(f"[delta] Parsing notebook: {path}")
                text = extract_notebook_cells(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "notebook", "path": path}))
            elif ext == ".py":
                logger.info(f"[delta] Parsing Python code (AST): {path}")
                text = extract_code_with_ast(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "python", "path": path}))
            elif ext in [".java", ".cpp", ".js", ".c"]:
                logger.info(f"[delta] Parsing {ext} code with Tree-sitter: {path}")
                lang = ext.strip(".")
                text = extract_code_with_treesitter(path, language=lang)
                if text:
                    docs.append(Document(text=text, metadata={"type": "code", "lang": lang, "path": path}))
            else:
                logger.info(f"[delta] Fallback reader for: {path}")
                fallback = SimpleDirectoryReader(input_files=[path])
                loaded = fallback.load_data()
                docs.extend(loaded)
        except Exception as e:
            logger.error(f"Failed to parse {path}: {e}")
    return docs

# New helpers: metadata load/save
def _meta_path(index_dir):
    return os.path.join(index_dir, EMBEDDING_META_FILENAME)

def load_embedding_metadata(index_dir):
    """Return a set of absolute file paths that were already embedded."""
    meta_file = _meta_path(index_dir)
    try:
        with open(meta_file, "r", encoding="utf-8") as f:
            data = json.load(f)
        return set(data.get("files", []))
    except Exception:
        return set()

def save_embedding_metadata(index_dir, files_set):
    """Persist the set of absolute file paths that have embeddings."""
    os.makedirs(index_dir, exist_ok=True)
    meta_file = _meta_path(index_dir)
    try:
        with open(meta_file, "w", encoding="utf-8") as f:
            json.dump({"files": sorted(list(files_set))}, f, indent=2)
    except Exception as e:
        logger.warning(f"Could not write embedding metadata file: {e}")

# ----------------------------- INDEX MANAGEMENT -----------------------------


def _index_has_expected_files(index_dir: str, expected_files=EXPECTED_PERSIST_FILES) -> bool:
    """Return True if all expected files exist in the index directory."""
    try:
        if not os.path.isdir(index_dir):
            return False
        contents = set(os.listdir(index_dir))
        missing = [f for f in expected_files if f not in contents]
        if missing:
            logger.debug(f"Missing expected index files: {missing}")
            return False
        return True
    except Exception as e:
        logger.warning(f"Could not verify index directory contents: {e}")
        return False


def get_or_create_index(index_dir: str = INDEX_DIR, data_dir: str = DATA_DIR, force_rebuild: bool = False):
    """Reuse existing vector index if available, otherwise rebuild.
    Verifies persistence files and recovers from corrupted index directories.
    """
    index_dir_abs = os.path.abspath(index_dir)
    logger.info(f"Index path: {index_dir_abs}")

    if force_rebuild:
        logger.warning("FORCE_REBUILD requested: will rebuild the index from documents.")
        # attempt to remove existing dir to ensure clean rebuild; if removal fails, we'll continue
        try:
            if os.path.isdir(index_dir_abs):
                shutil.rmtree(index_dir_abs)
                logger.info("Removed existing index directory due to force rebuild.")
        except Exception as e:
            logger.warning(f"Could not remove existing index directory: {e}")

    # Try to load existing index if persistence looks complete
    if _index_has_expected_files(index_dir_abs) and not force_rebuild:
        logger.info(f"🟢 Found existing index directory with expected files at {index_dir_abs}, attempting to load...")
        try:
            storage_context = StorageContext.from_defaults(persist_dir=index_dir_abs)
            index = load_index_from_storage(storage_context)
            logger.info("✅ Loaded existing index successfully.")

            # --- NEW: check embedding metadata for new files to embed ---
            try:
                # collect absolute paths of current data files
                current_files = set()
                for root, _, files in os.walk(data_dir):
                    for fname in files:
                        current_files.add(os.path.abspath(os.path.join(root, fname)))

                recorded = load_embedding_metadata(index_dir_abs)

                if not recorded:
                    # metadata missing: create metadata from current files (we assume existing index covers them)
                    save_embedding_metadata(index_dir_abs, current_files)
                    logger.info("Initialized embedding metadata from current files.")
                else:
                    new_files = current_files - recorded
                    if new_files:
                        logger.info(f"Found {len(new_files)} new file(s) to embed.")
                        new_docs = load_documents_for_paths(list(new_files))
                        if new_docs:
                            try:
                                # try incremental insertion using available API
                                if hasattr(index, "insert_documents"):
                                    index.insert_documents(new_docs)
                                elif hasattr(index, "add_documents"):
                                    index.add_documents(new_docs)
                                else:
                                    # fallback: rebuild entire index from all documents (safer fallback)
                                    logger.warning("Index doesn't support incremental insert; rebuilding full index to add new files.")
                                    all_docs = load_multimodal_documents(data_dir)
                                    parser = SimpleNodeParser.from_defaults(chunk_size=1000, chunk_overlap=200)
                                    nodes = parser.get_nodes_from_documents(all_docs)
                                    index = VectorStoreIndex(nodes)

                                # persist and update metadata
                                try:
                                    index.storage_context.persist(persist_dir=index_dir_abs)
                                except Exception as perr:
                                    logger.warning(f"Could not persist index after adding new docs: {perr}")
                                save_embedding_metadata(index_dir_abs, recorded.union(new_files))
                                logger.info("✅ Added embeddings for new files and updated metadata.")
                            except Exception as e:
                                logger.error(f"Failed to add new documents to index: {e}")
                        else:
                            logger.info("No documents were created from new files (skipping).")
            except Exception as meta_err:
                logger.warning(f"Embedding metadata check failed: {meta_err}")

            return index
        except Exception as e:
            logger.error(f"Failed to load existing index (will attempt to rebuild). Error: {e}")
            # move corrupt folder out of the way so we can create a fresh one
            corrupt_backup = index_dir_abs + ".corrupt"
            try:
                if os.path.exists(corrupt_backup):
                    shutil.rmtree(corrupt_backup)
                os.rename(index_dir_abs, corrupt_backup)
                logger.warning(f"Moved corrupt index folder to: {corrupt_backup}")
            except Exception as rename_err:
                logger.warning(f"Could not move corrupt index folder: {rename_err}. Will attempt to overwrite.")

    # Build a new index
    logger.warning("⚠️  Building a new index from documents...")
    documents = load_multimodal_documents(data_dir)
    if not documents:
        logger.error("❌ No documents found to index. Check DATA_DIR and files.")
        raise SystemExit(1)

    parser = SimpleNodeParser.from_defaults(chunk_size=1000, chunk_overlap=200)
    nodes = parser.get_nodes_from_documents(documents)
    index = VectorStoreIndex(nodes)

    # ensure directory exists then persist
    os.makedirs(index_dir_abs, exist_ok=True)
    try:
        index.storage_context.persist(persist_dir=index_dir_abs)
    except Exception as e:
        logger.error(f"Failed to persist index to {index_dir_abs}: {e}")
        raise

    # --- NEW: save metadata of all current files as having been embedded ---
    try:
        current_files = set()
        for root, _, files in os.walk(data_dir):
            for fname in files:
                current_files.add(os.path.abspath(os.path.join(root, fname)))
        save_embedding_metadata(index_dir_abs, current_files)
        logger.info("Saved embedding metadata for newly built index.")
    except Exception as e:
        logger.warning(f"Could not save embedding metadata after building index: {e}")

    logger.info(f"✅ Index built and persisted to {index_dir_abs}")
    return index

