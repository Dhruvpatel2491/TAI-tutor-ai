"""
Vector Store Generator for TAI Tutor AI.

This module handles document parsing, embedding generation, and vector index
management for RAG (Retrieval-Augmented Generation) functionality.

Supports:
- PDF extraction (PyMuPDF)
- PowerPoint extraction (python-pptx)
- Jupyter Notebook parsing
- Source code parsing (Python AST, Tree-sitter for multiple languages)
- Incremental index updates

Storage:
- Index persisted to INDEX_DIR (default: ./vector_index_store)
- Embedding metadata tracked in embeddings_meta.json
"""

import os

# Suppress TensorFlow warnings if present
os.environ["TF_ENABLE_ONEDNN_OPTS"] = "0"

import ast
import json
import logging
import shutil
from pathlib import Path
from typing import Any, Dict, List, Optional, Set

from pydantic import BaseModel

# Import configuration with fallback for running as script
try:
    from config import (
        INDEX_DIR,
        DATA_DIR,
        OLLAMA_LLM,
        OLLAMA_EMBED,
        DEFAULT_TEMPERATURE,
        DEFAULT_MAX_TOKENS,
        DEFAULT_TIMEOUT,
        EMBEDDINGS_DIR,
    )
except ImportError:
    from config import (
        INDEX_DIR,
        DATA_DIR,
        OLLAMA_LLM,
        OLLAMA_EMBED,
        DEFAULT_TEMPERATURE,
        DEFAULT_MAX_TOKENS,
        DEFAULT_TIMEOUT,
        EMBEDDINGS_DIR,
    )

logger = logging.getLogger("backend.rag.vector_store")

# =============================================================================
# External Library Imports (Optional)
# =============================================================================

# Try to import document processing libraries
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
    logger.warning("PyMuPDF (fitz) not installed. PDF extraction disabled.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logger.warning("python-pptx not installed. PPTX extraction disabled.")

try:
    from tree_sitter import Language, Parser
except ImportError:
    Language = None
    Parser = None
    logger.warning("tree-sitter not installed. Code parsing fallback to raw text.")

try:
    import nbformat
except ImportError:
    nbformat = None
    logger.warning("nbformat not installed. Notebook extraction disabled.")

# Try to import LlamaIndex components
try:
    from llama_index.core import (
        SimpleDirectoryReader,
        VectorStoreIndex,
        StorageContext,
        load_index_from_storage,
    )
    from llama_index.core.node_parser import SimpleNodeParser, SentenceSplitter
    from llama_index.core.schema import Document
    from llama_index.core.settings import Settings
    from llama_index.llms.ollama import Ollama
    from llama_index.embeddings.ollama import OllamaEmbedding
    LLAMA_INDEX_AVAILABLE = True
except ImportError:
    LLAMA_INDEX_AVAILABLE = False
    logger.warning("llama_index not fully installed. Vector store functionality limited.")


# =============================================================================
# Configuration
# =============================================================================

EXPECTED_PERSIST_FILES = ["docstore.json"]
EMBEDDING_META_FILENAME = "embeddings_meta.json"


# =============================================================================
# Document Parsers
# =============================================================================

def extract_text_from_pdf(filepath: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    if fitz is None:
        logger.error("PyMuPDF not available for PDF extraction")
        return ""
    try:
        with fitz.open(filepath) as doc:
            text = [page.get_text("text") for page in doc]
        return "\n".join(text)
    except Exception as e:
        logger.error(f"PDF extraction failed for {filepath}: {e}")
        return ""


def extract_text_from_pptx(filepath: str) -> str:
    """Extract slide-level text using python-pptx."""
    if Presentation is None:
        logger.error("python-pptx not available for PPTX extraction")
        return ""
    try:
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())
            if content:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(content))
        return "\n\n".join(slides_text)
    except Exception as e:
        logger.error(f"PPTX extraction failed for {filepath}: {e}")
        return ""


def extract_code_with_treesitter(filepath: str, language: str = "python") -> str:
    """Parse source code using Tree-sitter (multi-language)."""
    if Language is None or Parser is None:
        # Fallback to raw code
        try:
            return Path(filepath).read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""
    
    lang_so = "build/my-languages.so"
    
    # Build library if missing
    if not os.path.exists(lang_so):
        logger.info("Building Tree-sitter language library (this may take a moment)...")
        os.makedirs("vendor", exist_ok=True)
        try:
            Language.build_library(
                lang_so,
                [
                    "vendor/tree-sitter-python",
                    "vendor/tree-sitter-javascript",
                    "vendor/tree-sitter-cpp",
                    "vendor/tree-sitter-java",
                ],
            )
        except Exception as e:
            logger.warning(f"Failed to build Tree-sitter library: {e}")
            return Path(filepath).read_text(encoding="utf-8", errors="ignore")

    try:
        lang = Language(lang_so, language)
    except Exception as e:
        logger.error(f"Tree-sitter Language init failed for '{language}': {e}")
        return Path(filepath).read_text(encoding="utf-8", errors="ignore")

    parser = Parser()
    parser.set_language(lang)

    code = Path(filepath).read_text(encoding="utf-8", errors="ignore")
    try:
        _ = parser.parse(bytes(code, "utf8"))
    except Exception as e:
        logger.warning(f"Tree-sitter parse warning for {filepath}: {e}")
    return code


class NotebookCell(BaseModel):
    """Represents a Jupyter notebook cell."""
    cell_type: str
    source: str
    metadata: dict = {}


def extract_notebook_cells(filepath: str) -> str:
    """Extract notebook code and markdown cells as JSON."""
    if nbformat is None:
        logger.error("nbformat not available for notebook extraction")
        return ""
    try:
        nb = nbformat.read(filepath, as_version=4)
        safe_cells = []
        for c in nb.cells:
            nc = NotebookCell(
                cell_type=c.cell_type,
                source="".join(c.source) if isinstance(c.source, (list, tuple)) else str(c.source),
                metadata=getattr(c, "metadata", {}),
            )
            try:
                safe_cells.append(nc.model_dump())
            except Exception:
                safe_cells.append(getattr(nc, "__dict__", {}))
        return json.dumps(safe_cells, indent=2)
    except Exception as e:
        logger.error(f"Notebook parse failed for {filepath}: {e}")
        return ""


def extract_code_with_ast(filepath: str) -> str:
    """Analyze Python code structure with AST (functions/classes)."""
    try:
        src = Path(filepath).read_text(encoding="utf-8")
        tree = ast.parse(src)
        funcs = [n.name for n in ast.walk(tree) if isinstance(n, ast.FunctionDef)]
        classes = [n.name for n in ast.walk(tree) if isinstance(n, ast.ClassDef)]
        summary = f"Functions: {funcs}\nClasses: {classes}\n\n{src}"
        return summary
    except Exception as e:
        logger.error(f"AST parse failed for {filepath}: {e}")
        return ""


# =============================================================================
# Document Loading
# =============================================================================

def load_multimodal_documents(directory: str) -> List[Any]:
    """
    Load PDFs, PPTXs, Notebooks, and code files from a directory.
    
    Returns list of LlamaIndex Document objects.
    """
    if not LLAMA_INDEX_AVAILABLE:
        logger.error("LlamaIndex not available for document loading")
        return []
    
    docs = []
    
    if not os.path.exists(directory):
        logger.warning(f"Data directory does not exist: {directory}")
        return docs

    for root, _, files in os.walk(directory):
        for file in files:
            path = os.path.join(root, file)
            ext = Path(file).suffix.lower()

            try:
                if ext == ".pdf":
                    logger.info(f"Parsing PDF: {path}")
                    text = extract_text_from_pdf(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "pdf", "path": path}))

                elif ext == ".pptx":
                    logger.info(f"Parsing PPTX: {path}")
                    text = extract_text_from_pptx(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "pptx", "path": path}))

                elif ext == ".ipynb":
                    logger.info(f"Parsing notebook: {path}")
                    text = extract_notebook_cells(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "notebook", "path": path}))

                elif ext == ".py":
                    logger.info(f"Parsing Python code (AST): {path}")
                    text = extract_code_with_ast(path)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "python", "path": path}))

                elif ext in [".java", ".cpp", ".js", ".c"]:
                    logger.info(f"Parsing {ext} code with Tree-sitter: {path}")
                    lang = ext.strip(".")
                    text = extract_code_with_treesitter(path, language=lang)
                    if text:
                        docs.append(Document(text=text, metadata={"type": "code", "lang": lang, "path": path}))

                else:
                    logger.info(f"Fallback reader for: {path}")
                    try:
                        fallback = SimpleDirectoryReader(input_files=[path])
                        loaded = fallback.load_data()
                        docs.extend(loaded)
                    except Exception:
                        logger.debug(f"Could not load {path} with fallback reader")

            except Exception as e:
                logger.error(f"Failed to parse {path}: {e}")

    return docs


def load_documents_for_paths(paths: List[str]) -> List[Any]:
    """
    Load specific file paths into LlamaIndex Documents.
    
    Used for incremental updates when new files are detected.
    """
    if not LLAMA_INDEX_AVAILABLE:
        logger.error("LlamaIndex not available for document loading")
        return []
    
    docs = []
    for path in paths:
        try:
            if not os.path.exists(path):
                logger.warning(f"Path not found while building selective docs: {path}")
                continue
            ext = Path(path).suffix.lower()
            
            if ext == ".pdf":
                logger.info(f"[delta] Parsing PDF: {path}")
                text = extract_text_from_pdf(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "pdf", "path": path}))
                    
            elif ext == ".pptx":
                logger.info(f"[delta] Parsing PPTX: {path}")
                text = extract_text_from_pptx(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "pptx", "path": path}))
                    
            elif ext == ".ipynb":
                logger.info(f"[delta] Parsing notebook: {path}")
                text = extract_notebook_cells(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "notebook", "path": path}))
                    
            elif ext == ".py":
                logger.info(f"[delta] Parsing Python code (AST): {path}")
                text = extract_code_with_ast(path)
                if text:
                    docs.append(Document(text=text, metadata={"type": "python", "path": path}))
                    
            elif ext in [".java", ".cpp", ".js", ".c"]:
                logger.info(f"[delta] Parsing {ext} code with Tree-sitter: {path}")
                lang = ext.strip(".")
                text = extract_code_with_treesitter(path, language=lang)
                if text:
                    docs.append(Document(text=text, metadata={"type": "code", "lang": lang, "path": path}))
                    
            else:
                logger.info(f"[delta] Fallback reader for: {path}")
                try:
                    fallback = SimpleDirectoryReader(input_files=[path])
                    loaded = fallback.load_data()
                    docs.extend(loaded)
                except Exception:
                    logger.debug(f"Could not load {path} with fallback reader")
                    
        except Exception as e:
            logger.error(f"Failed to parse {path}: {e}")
    return docs


# =============================================================================
# Embedding Metadata
# =============================================================================

def _meta_path(index_dir: str) -> str:
    """Get path to embedding metadata file."""
    return os.path.join(index_dir, EMBEDDING_META_FILENAME)


def load_embedding_metadata(index_dir: str) -> Set[str]:
    """Return set of absolute file paths that were already embedded."""
    meta_file = _meta_path(index_dir)
    try:
        with open(meta_file, "r", encoding="utf-8") as f:
            data = json.load(f)
        return set(data.get("files", []))
    except Exception:
        return set()


def save_embedding_metadata(index_dir: str, files_set: Set[str]) -> None:
    """Persist the set of absolute file paths that have embeddings."""
    os.makedirs(index_dir, exist_ok=True)
    meta_file = _meta_path(index_dir)
    try:
        with open(meta_file, "w", encoding="utf-8") as f:
            json.dump({"files": sorted(list(files_set))}, f, indent=2)
    except Exception as e:
        logger.warning(f"Could not write embedding metadata file: {e}")


# =============================================================================
# Index Management
# =============================================================================

def _index_has_expected_files(index_dir: str, expected_files: List[str] = None) -> bool:
    """Return True if all expected files exist in the index directory."""
    if expected_files is None:
        expected_files = EXPECTED_PERSIST_FILES
    try:
        if not os.path.isdir(index_dir):
            return False
        contents = set(os.listdir(index_dir))
        missing = [f for f in expected_files if f not in contents]
        if missing:
            logger.debug(f"Missing expected index files: {missing}")
            return False
        return True
    except Exception as e:
        logger.warning(f"Could not verify index directory contents: {e}")
        return False


def _read_embedding_model_from_folder(embeddings_dir: str, fallback: str = "") -> str:
    """Read embedding model name from embeddings folder configuration."""
    try:
        p_txt = os.path.join(embeddings_dir, "model.txt")
        if os.path.exists(p_txt):
            with open(p_txt, "r", encoding="utf-8") as f:
                val = f.read().strip()
                if val:
                    return val
        p_json = os.path.join(embeddings_dir, "config.json")
        if os.path.exists(p_json):
            with open(p_json, "r", encoding="utf-8") as f:
                cfg = json.load(f)
            if isinstance(cfg, dict) and "model" in cfg and cfg["model"]:
                return str(cfg["model"])
    except Exception:
        logger.debug("Could not read embedding model from folder; using fallback.")
    return fallback or OLLAMA_EMBED


def init_models() -> None:
    """Initialize LlamaIndex Settings for embedding and LLM models."""
    if not LLAMA_INDEX_AVAILABLE:
        raise RuntimeError("LlamaIndex not available for model initialization")
    try:
        embed_model_name = _read_embedding_model_from_folder(EMBEDDINGS_DIR, fallback=OLLAMA_EMBED)
        Settings.embed_model = OllamaEmbedding(
            model_name=embed_model_name,
            max_tokens=DEFAULT_MAX_TOKENS,
            request_timeout=DEFAULT_TIMEOUT
        )
        Settings.llm = Ollama(
            model=OLLAMA_LLM,
            temperature=DEFAULT_TEMPERATURE,
            max_tokens=DEFAULT_MAX_TOKENS,
            request_timeout=DEFAULT_TIMEOUT
        )
        logger.info(f"Initialized Ollama LLM='{OLLAMA_LLM}' embed='{embed_model_name}'")
    except Exception as e:
        logger.error(f"Could not initialize Ollama models: {e}")
        raise


def get_or_create_index(
    index_dir: str = None,
    data_dir: str = None,
    force_rebuild: bool = False,
    indexing: Optional[Dict[str, Any]] = None
) -> Any:
    """
    Reuse existing vector index if available, otherwise rebuild.
    
    Verifies persistence files and recovers from corrupted index directories.
    Supports build-time indexing parameters for parser selection.
    
    Args:
        index_dir: Directory for persisted index
        data_dir: Directory containing source documents
        force_rebuild: Force complete rebuild of index
        indexing: Build-time parameters (chunk_size, chunk_overlap, parser, etc.)
    
    Returns:
        VectorStoreIndex instance
    """
    if not LLAMA_INDEX_AVAILABLE:
        raise RuntimeError("LlamaIndex not available for index creation")
    
    index_dir = index_dir or INDEX_DIR
    data_dir = data_dir or DATA_DIR
    index_dir_abs = os.path.abspath(index_dir)
    logger.info(f"Index path: {index_dir_abs}")

    if force_rebuild:
        logger.warning("FORCE_REBUILD requested: will rebuild the index from documents.")
        try:
            if os.path.isdir(index_dir_abs):
                shutil.rmtree(index_dir_abs)
                logger.info("Removed existing index directory due to force rebuild.")
        except Exception as e:
            logger.warning(f"Could not remove existing index directory: {e}")

    # Try to load existing index
    if _index_has_expected_files(index_dir_abs) and not force_rebuild:
        logger.info(f"🟢 Found existing index at {index_dir_abs}, attempting to load...")
        try:
            storage_context = StorageContext.from_defaults(persist_dir=index_dir_abs)
            index = load_index_from_storage(storage_context)
            logger.info("✅ Loaded existing index successfully.")

            # Check for new files to embed
            try:
                current_files = set()
                for root, _, files in os.walk(data_dir):
                    for fname in files:
                        current_files.add(os.path.abspath(os.path.join(root, fname)))

                recorded = load_embedding_metadata(index_dir_abs)

                if not recorded:
                    save_embedding_metadata(index_dir_abs, current_files)
                    logger.info("Initialized embedding metadata from current files.")
                else:
                    new_files = current_files - recorded
                    if new_files:
                        logger.info(f"Found {len(new_files)} new file(s) to embed.")
                        new_docs = load_documents_for_paths(list(new_files))
                        if new_docs:
                            try:
                                if hasattr(index, "insert_documents"):
                                    index.insert_documents(new_docs)
                                elif hasattr(index, "add_documents"):
                                    index.add_documents(new_docs)
                                else:
                                    logger.warning("Index doesn't support incremental insert; rebuilding.")
                                    all_docs = load_multimodal_documents(data_dir)
                                    parser = SimpleNodeParser.from_defaults(
                                        chunk_size=1000,
                                        chunk_overlap=200,
                                        separator="\n\n",
                                        include_metadata=True
                                    )
                                    nodes = parser.get_nodes_from_documents(all_docs)
                                    index = VectorStoreIndex(nodes)

                                try:
                                    index.storage_context.persist(persist_dir=index_dir_abs)
                                except Exception as perr:
                                    logger.warning(f"Could not persist index after adding new docs: {perr}")
                                save_embedding_metadata(index_dir_abs, recorded.union(new_files))
                                logger.info("✅ Added embeddings for new files.")
                            except Exception as e:
                                logger.error(f"Failed to add new documents to index: {e}")
            except Exception as meta_err:
                logger.warning(f"Embedding metadata check failed: {meta_err}")

            return index
            
        except Exception as e:
            logger.error(f"Failed to load existing index (will rebuild). Error: {e}")
            corrupt_backup = index_dir_abs + ".corrupt"
            try:
                if os.path.exists(corrupt_backup):
                    shutil.rmtree(corrupt_backup)
                os.rename(index_dir_abs, corrupt_backup)
                logger.warning(f"Moved corrupt index folder to: {corrupt_backup}")
            except Exception as rename_err:
                logger.warning(f"Could not move corrupt index folder: {rename_err}")

    # Build new index
    logger.warning("⚠️ Building a new index from documents...")
    documents = load_multimodal_documents(data_dir)
    if not documents:
        logger.error("❌ No documents found to index. Check DATA_DIR and files.")
        raise SystemExit(1)

    # Parse indexing parameters
    if indexing is None:
        indexing = {}

    parser_type = indexing.get("parser") or indexing.get("splitter") or "simple"
    chunk_size = int(indexing.get("chunk_size", 1000))
    chunk_overlap = int(indexing.get("chunk_overlap", 200))
    separator = indexing.get("separator", "\n\n")
    include_metadata = bool(indexing.get("include_metadata", True))

    # Select parser
    if parser_type.lower() in ["sentence", "sentencesplitter"]:
        parser = SentenceSplitter.from_defaults(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separator=separator,
            include_metadata=include_metadata
        )
    else:
        parser = SimpleNodeParser.from_defaults(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separator=separator,
            include_metadata=include_metadata
        )

    nodes = parser.get_nodes_from_documents(documents)
    index = VectorStoreIndex(nodes)

    # Persist index
    os.makedirs(index_dir_abs, exist_ok=True)
    try:
        index.storage_context.persist(persist_dir=index_dir_abs)
    except Exception as e:
        logger.error(f"Failed to persist index to {index_dir_abs}: {e}")
        raise

    # Save embedding metadata
    try:
        current_files = set()
        for root, _, files in os.walk(data_dir):
            for fname in files:
                current_files.add(os.path.abspath(os.path.join(root, fname)))
        save_embedding_metadata(index_dir_abs, current_files)
        logger.info("Saved embedding metadata for newly built index.")
    except Exception as e:
        logger.warning(f"Could not save embedding metadata after building index: {e}")

    logger.info(f"✅ Index built and persisted to {index_dir_abs}")
    return index
